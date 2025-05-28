
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd
import tempfile
import os

st.title("PCA Commentary Generator with Summary Charts")

pptx_file = st.file_uploader("Upload PowerPoint Template", type=["pptx"])
excel_file = st.file_uploader("Upload Excel Dataset", type=["xlsx"])
slide_indices_input = st.text_input("Slide indices (comma-separated)", value="5,6,7")
generate_button = st.button("Generate Commentary & Summary Charts")

def generate_commentary_from_excel(prs, df, slide_indices):
    data_rows = df[df[df.columns[0]].isna() & df[df.columns[5]].notna()]
    data_rows = data_rows.reset_index(drop=True)

    def generate_commentary_row(row):
        commentary = []
        try:
            planned_impressions = float(row[df.columns[5]])
            actual_impressions = float(row[df.columns[6]])
            planned_cpm = float(row[df.columns[7]])
            actual_cpm = float(row[df.columns[8]]) if pd.notna(row[df.columns[8]]) else planned_cpm
            ctr = float(row[df.columns[9]])

            if planned_impressions > 0:
                imp_diff = ((actual_impressions - planned_impressions) / planned_impressions) * 100
                direction = "higher" if imp_diff > 0 else "lower"
                commentary.append(f"Impressions were {abs(imp_diff):.1f}% {direction} than planned.")

            if planned_cpm > 0:
                cpm_diff = ((actual_cpm - planned_cpm) / planned_cpm) * 100
                direction = "higher" if cpm_diff > 0 else "lower"
                commentary.append(f"CPM was {abs(cpm_diff):.1f}% {direction} than planned.")

            if ctr >= 0.07:
                commentary.append("CTR met or exceeded the 0.07% benchmark.")
            else:
                commentary.append("CTR was below the 0.07% benchmark.")

        except Exception:
            pass

        return " ".join(commentary)

    valid_indices = [i for i in slide_indices if i < len(prs.slides)]
    for i, slide_index in enumerate(valid_indices):
        row = data_rows.iloc[i]
        commentary_text = generate_commentary_row(row)
        slide = prs.slides[slide_index]

        if commentary_text:
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(8.5), Inches(1))
            text_frame = textbox.text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = commentary_text
            font = run.font
            font.size = Pt(14)
            font.name = 'Arial'
            font.color.rgb = RGBColor(0, 0, 0)

    # Add summary slide
    summary_slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Pie Chart: Total Digital Media Cost by Placement
    cost_data = df[df[df.columns[0]].isna() & df[df.columns[4]].notna()]
    cost_data = cost_data.reset_index(drop=True)
    pie_data = CategoryChartData()
    pie_data.categories = cost_data[df.columns[2]].fillna("Unknown Placement")
    pie_data.add_series('Cost', [float(val) if pd.notna(val) else 0 for val in cost_data[df.columns[4]]])
    summary_slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.5), Inches(0.5), Inches(4.5), Inches(3.5), pie_data)

    # Bar Chart: CPM comparison (Planned vs Actual)
    bar_data = CategoryChartData()
    bar_data.categories = cost_data[df.columns[2]].fillna("Unknown Placement")
    bar_data.add_series('Planned CPM', [float(val) if pd.notna(val) else 0 for val in cost_data[df.columns[7]]])
    bar_data.add_series('Actual CPM', [float(val) if pd.notna(val) else 0 for val in cost_data[df.columns[8]]])
    summary_slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5.2), Inches(0.5), Inches(4.5), Inches(3.5), bar_data)

    return prs

if generate_button and pptx_file and excel_file:
    slide_indices = [int(x.strip()) for x in slide_indices_input.split(",") if x.strip().isdigit()]
    df = pd.read_excel(excel_file)
    prs = Presentation(pptx_file)

    prs = generate_commentary_from_excel(prs, df, slide_indices)

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        prs.save(tmp.name)
        tmp_path = tmp.name

    with open(tmp_path, "rb") as file:
        st.download_button(
            label="Download Presentation with Commentary & Summary Charts",
            data=file,
            file_name="Updated_PCA_with_Summary_Charts.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    os.unlink(tmp_path)
